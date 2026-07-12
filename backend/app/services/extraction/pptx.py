"""PowerPoint (.pptx) extractor via python-pptx.

Each slide becomes one page, so answers cite the slide number. Text is pulled
from every shape's text frame and from table cells, in slide order.
"""

from __future__ import annotations

from pathlib import Path

from app.core.exceptions import ExtractionError
from app.models.document import Document
from app.services.extraction.base import BaseExtractor


class PptxExtractor(BaseExtractor):
    extensions = (".pptx",)
    mime_types = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    def extract(self, path: Path) -> Document:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - env dependent
            raise ExtractionError(
                "python-pptx is required for PPTX extraction. "
                "Install with `pip install python-pptx`."
            ) from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ExtractionError(f"failed to parse PPTX {path.name}: {exc}") from exc

        slide_texts: list[str] = []
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            slide_texts.append("\n".join(parts))

        return self.assemble(
            filename=path.name,
            mime_type=self.mime_types[0],
            page_texts=slide_texts,
            metadata={"extractor": "python-pptx"},
        )
